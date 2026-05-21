"""Tests for pptx_create tool."""

import asyncio
from pathlib import Path

import pytest

from src.config import settings
from src.tools.pptx_create import pptx_create_handler


class FakeServer:
    def __init__(self):
        self.tools = {}

    def tool(self, *args, **kwargs):
        def decorator(fn):
            self.tools[fn.__name__] = fn
            return fn
        return decorator


def _get_pptx_fn():
    server = FakeServer()
    pptx_create_handler(server)
    return server.tools["pptx_create"]


def _run(fn, **kwargs):
    return asyncio.run(fn(**kwargs))


@pytest.fixture
def output_dir(tmp_path):
    """Use a temporary output directory."""
    settings.FILE_OUTPUT_DIR = str(tmp_path)
    return tmp_path


class TestPptxCreate:
    """Tests for the pptx_create tool."""

    def test_basic_presentation(self, output_dir):
        """Create a simple presentation with a single title slide."""
        fn = _get_pptx_fn()
        result = _run(fn, filename="test_basic.pptx", slides=[{
            "layout": "title",
            "title": "Hello World",
        }])

        assert len(result) == 2  # embedded resource + text content
        assert result[0].type == "resource"
        assert result[1].type == "text"

        # Verify file was created
        file_path = output_dir / "test_basic.pptx"
        assert file_path.exists()
        assert file_path.suffix == ".pptx"

    def test_multiple_slides(self, output_dir):
        """Create a presentation with multiple slides."""
        fn = _get_pptx_fn()
        result = _run(fn, filename="test_multi.pptx", slides=[
            {
                "layout": "title",
                "title": "Title Slide",
            },
            {
                "layout": "section_header",
                "title": "Second Slide",
            },
            {
                "layout": "title_content",
                "title": "Content Slide",
                "body": ["Point 1", "Point 2", "Point 3"],
            },
        ])

        assert len(result) == 2
        assert "3" in result[1].text  # 3 slides

        file_path = output_dir / "test_multi.pptx"
        assert file_path.exists()

    def test_body_text(self, output_dir):
        """Create a slide with body text."""
        fn = _get_pptx_fn()
        result = _run(fn, filename="test_body.pptx", slides=[{
            "layout": "title_content",
            "title": "Agenda",
            "body": ["Introduction", "Discussion", "Q&A"],
        }])

        assert result[0].type == "resource"
        file_path = output_dir / "test_body.pptx"
        assert file_path.exists()

    def test_bullet_text(self, output_dir):
        """Create a slide with bullet points."""
        fn = _get_pptx_fn()
        result = _run(fn, filename="test_bullet.pptx", slides=[{
            "layout": "title_content",
            "title": "Key Points",
            "body": ["First point", "Second point"],
            "bullet": True,
        }])

        assert result[0].type == "resource"
        file_path = output_dir / "test_bullet.pptx"
        assert file_path.exists()

    def test_notes(self, output_dir):
        """Create a slide with speaker notes."""
        fn = _get_pptx_fn()
        result = _run(fn, filename="test_notes.pptx", slides=[{
            "layout": "section_header",
            "title": "Notes Slide",
            "notes": "Remember to mention the budget",
        }])

        assert result[0].type == "resource"
        file_path = output_dir / "test_notes.pptx"
        assert file_path.exists()

    def test_theme_customization(self, output_dir):
        """Create a presentation with dark theme."""
        fn = _get_pptx_fn()
        result = _run(fn, filename="test_theme.pptx",
            theme="dark",
            slides=[{
                "layout": "title",
                "title": "Dark Theme",
            }],
        )

        assert result[0].type == "resource"
        file_path = output_dir / "test_theme.pptx"
        assert file_path.exists()

    def test_invalid_layout(self, output_dir):
        """Error is returned for invalid layout name."""
        fn = _get_pptx_fn()
        result = _run(fn, filename="test_invalid.pptx", slides=[{
            "layout": "nonexistent_layout",
            "title": "Bad",
        }])

        # On error, returns only TextContent with error message
        assert len(result) == 1
        assert result[0].type == "text"
        assert "Error" in result[0].text
        assert "nonexistent_layout" in result[0].text

    def test_no_slides(self, output_dir):
        """Error is returned when no slides are provided."""
        fn = _get_pptx_fn()
        result = _run(fn, filename="test_empty.pptx", slides=[])

        assert len(result) == 1
        assert result[0].type == "text"
        assert "Error" in result[0].text

    def test_handler_integration(self, output_dir):
        """Verify the handler function works with an MCP server."""
        from mcp.server.fastmcp import FastMCP

        server = FastMCP(name="test")
        pptx_create_handler(server)

        # Verify the tool is registered
        tool_names = [t.name for t in server._tool_manager.list_tools()]
        assert "pptx_create" in tool_names

    def test_output_file_is_valid_pptx(self, output_dir):
        """Verify the output file can be opened by python-pptx."""
        fn = _get_pptx_fn()
        _run(fn, filename="test_valid.pptx", slides=[{
            "layout": "title_content",
            "title": "Valid Slide",
            "body": ["Some content"],
        }])

        # Re-open the file to verify it's a valid pptx
        from pptx import Presentation
        file_path = output_dir / "test_valid.pptx"
        prs = Presentation(str(file_path))

        assert len(prs.slides) == 1
        # Our tool adds title via a new textbox (not the layout placeholder)
        # so search all shapes for the title text
        slide = prs.slides[0]
        all_texts = [s.text_frame.text for s in slide.shapes]
        assert "Valid Slide" in all_texts
