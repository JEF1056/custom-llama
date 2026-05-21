"""PowerPoint (.pptx) creation tool for MCP server."""

import base64
import logging
from io import BytesIO
from pathlib import Path

from mcp.server import FastMCP
from mcp.types import EmbeddedResource, TextContent
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.config import settings

logger = logging.getLogger(__name__)

# Slide layout indices for the default PowerPoint template
_LAYOUT_MAP = {
    "title": 0,          # Title Slide
    "title_content": 1,  # Title and Content
    "section_header": 2, # Section Header
    "two_column": 3,     # Two Content
    "blank": 6,          # Blank
}

_VALID_LAYOUTS = set(_LAYOUT_MAP.keys())
_VALID_THEMES = {"default", "dark", "minimal"}
_VALID_IMAGE_POSITIONS = {"left", "right", "center"}

# Theme color definitions
_THEMES = {
    "default": {
        "bg": None,  # use slide master default
        "title_color": None,
        "body_color": None,
        "accent": RGBColor(0x44, 0x72, 0xC4),
    },
    "dark": {
        "bg": RGBColor(0x1E, 0x1E, 0x1E),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "body_color": RGBColor(0xCC, 0xCC, 0xCC),
        "accent": RGBColor(0x00, 0xB4, 0xD8),
    },
    "minimal": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x33, 0x33, 0x33),
        "body_color": RGBColor(0x55, 0x55, 0x55),
        "accent": RGBColor(0x66, 0x66, 0x66),
    },
}

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Margins and spacing
MARGIN = Inches(0.5)
CONTENT_TOP = Inches(1.8)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN * 2
COLUMN_GAP = Inches(0.4)
COLUMN_WIDTH = (CONTENT_WIDTH - COLUMN_GAP) / 2

# Table defaults
TABLE_TOP = Inches(2.0)
TABLE_HEIGHT = Inches(4.5)
TABLE_COL_MIN = Inches(1.0)


def _validate_filename(filename: str) -> str | None:
    """Validate filename for path traversal. Returns error message or None."""
    if not filename:
        return "Error: filename is required"
    if "/" in filename or "\\" in filename or ".." in filename:
        return "Error: filename must not contain path separators or '..'"
    if not filename.lower().endswith(".pptx"):
        return "Error: filename must end with .pptx"
    return None


def _set_slide_background(slide, theme: dict) -> None:
    """Set slide background color for the given theme."""
    if theme["bg"] is None:
        return
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme["bg"]


def _set_run_color(run, color: RGBColor | None) -> None:
    """Set font color on a run if specified."""
    if color is not None:
        run.font.color.rgb = color


def _add_title_text(slide, text: str, theme: dict,
                    left=None, top=None, width=None, height=None,
                    font_size=Pt(36), bold=True, alignment=PP_ALIGN.LEFT) -> None:
    """Add a title text box to the slide."""
    if left is None:
        left = MARGIN
    if top is None:
        top = Inches(0.5)
    if width is None:
        width = CONTENT_WIDTH
    if height is None:
        height = Inches(1.2)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    _set_run_color(run, theme["title_color"])


def _add_body_text(slide, content, theme: dict,
                    left=None, top=None, width=None, height=None,
                    font_size=Pt(18), bullet=False,
                    alignment=PP_ALIGN.LEFT) -> None:
    """Add body text (plain or bulleted) to the slide."""
    if left is None:
        left = MARGIN
    if top is None:
        top = CONTENT_TOP
    if width is None:
        width = CONTENT_WIDTH
    if height is None:
        height = Inches(4.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    items = content if isinstance(content, list) else [content]

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = alignment

        if bullet:
            p.level = 0
            p.text = ""
            run = p.add_run()
            run.text = f"\u2022 {item}"
        else:
            p.text = ""
            run = p.add_run()
            run.text = item

        run.font.size = font_size
        _set_run_color(run, theme["body_color"])


def _add_table(slide, table_def: dict, theme: dict,
               left=None, top=None, width=None, height=None) -> str | None:
    """Add a table to the slide. Returns error message or None."""
    headers = table_def.get("headers", [])
    rows = table_def.get("rows", [])

    if not headers:
        return "Error: table requires 'headers'"

    num_rows = 1 + len(rows)  # header + data rows
    num_cols = len(headers)

    if left is None:
        left = MARGIN
    if top is None:
        top = TABLE_TOP
    if width is None:
        width = CONTENT_WIDTH
    if height is None:
        height = TABLE_HEIGHT

    # Ensure minimum column width
    min_total = TABLE_COL_MIN * num_cols
    if width < min_total:
        width = min_total

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths evenly
    col_width = width // num_cols
    for i in range(num_cols):
        table.columns[i].width = col_width

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(header)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(14)
                if theme["title_color"]:
                    run.font.color.rgb = theme["title_color"]
                else:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.alignment = PP_ALIGN.CENTER
        # Header cell background
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = theme["accent"]

    # Data rows
    for r_idx, row in enumerate(rows, 1):
        for c_idx, value in enumerate(row):
            if c_idx >= num_cols:
                break
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value) if value is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    _set_run_color(run, theme["body_color"])

    return None


def _add_image(slide, image_path: Path, position: str, theme: dict) -> str | None:
    """Add an image to the slide. Returns error message or None."""
    if not image_path.exists():
        return f"Error: image file not found: {image_path.name}"

    img_width = Inches(5.0)
    img_height = Inches(4.0)

    if position == "left":
        left = MARGIN
        top = CONTENT_TOP
    elif position == "right":
        left = SLIDE_WIDTH - MARGIN - img_width
        top = CONTENT_TOP
    elif position == "center":
        left = (SLIDE_WIDTH - img_width) / 2
        top = (SLIDE_HEIGHT - img_height) / 2
    else:
        left = MARGIN
        top = CONTENT_TOP

    try:
        slide.shapes.add_picture(str(image_path), left, top, img_width, img_height)
    except Exception as e:
        return f"Error: failed to add image - {str(e)}"
    return None


def _add_notes(slide, notes_text: str) -> None:
    """Add speaker notes to the slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def _set_title(slide, text: str, theme: dict, font_size=Pt(36), bold=True,
               alignment=PP_ALIGN.LEFT) -> None:
    """Set text on the slide's built-in title placeholder, or create a textbox."""
    title_shape = slide.shapes.title
    if title_shape and title_shape.has_text_frame:
        tf = title_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = alignment
        p.text = ""
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        _set_run_color(run, theme["title_color"])
    else:
        _add_title_text(slide, text, theme, font_size=font_size, bold=bold, alignment=alignment)


def _get_content(slide_def: dict):
    """Get content from slide_def, supporting both 'content' and 'body' keys."""
    content = slide_def.get("content")
    if content is None:
        content = slide_def.get("body")
    return content


def _get_bullet_flag(slide_def: dict, content) -> bool:
    """Determine if content should be rendered as bullets."""
    # Explicit 'bullet' key takes precedence
    if slide_def.get("bullet") is not None:
        return bool(slide_def["bullet"])
    # Otherwise, list content implies bullets
    return isinstance(content, list)


def _build_slide(prs: "Presentation", slide_def: dict, theme: dict,
                 file_output_dir: Path) -> str | None:
    """Build a single slide. Returns error message or None."""
    layout_name = slide_def.get("layout", "title_content")

    if layout_name not in _VALID_LAYOUTS:
        valid = ", ".join(sorted(_VALID_LAYOUTS))
        return f"Error: invalid layout '{layout_name}', must be {valid}"

    layout_idx = _LAYOUT_MAP[layout_name]

    # Handle layout index out of range by falling back to blank
    if layout_idx >= len(prs.slide_layouts):
        layout_idx = len(prs.slide_layouts) - 1

    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)

    # Apply theme background
    _set_slide_background(slide, theme)

    title = slide_def.get("title")
    subtitle = slide_def.get("subtitle")
    content = _get_content(slide_def)
    bullet = _get_bullet_flag(slide_def, content)
    notes = slide_def.get("notes")
    image = slide_def.get("image")
    image_position = slide_def.get("image_position", "center")
    table_def = slide_def.get("table")
    left_content = slide_def.get("left_content")
    right_content = slide_def.get("right_content")

    # --- Title slide layout ---
    if layout_name == "title":
        if title:
            _set_title(slide, title, theme, font_size=Pt(44),
                       alignment=PP_ALIGN.CENTER)
        if subtitle:
            _add_body_text(slide, subtitle, theme,
                          top=Inches(4.0), font_size=Pt(24),
                          alignment=PP_ALIGN.CENTER)

    # --- Section header layout ---
    elif layout_name == "section_header":
        if title:
            _set_title(slide, title, theme, font_size=Pt(40),
                       alignment=PP_ALIGN.CENTER)
        if subtitle:
            _add_body_text(slide, subtitle, theme,
                          top=Inches(4.2), font_size=Pt(22),
                          alignment=PP_ALIGN.CENTER)

    # --- Title + Content layout ---
    elif layout_name == "title_content":
        if title:
            _set_title(slide, title, theme, font_size=Pt(32))

        content_top = Inches(1.8)
        if table_def:
            err = _add_table(slide, table_def, theme, top=content_top)
            if err:
                return err
        elif image:
            img_path = file_output_dir / image
            err = _add_image(slide, img_path, image_position, theme)
            if err:
                return err
            if content:
                # Place text beside or below image
                if image_position in ("left", "right"):
                    text_left = (MARGIN if image_position == "right"
                                 else SLIDE_WIDTH - MARGIN - COLUMN_WIDTH)
                    _add_body_text(slide, content, theme,
                                  left=text_left, top=content_top,
                                  width=COLUMN_WIDTH)
                else:
                    _add_body_text(slide, content, theme,
                                  top=Inches(6.2), height=Inches(1.0))
        elif content:
            _add_body_text(slide, content, theme, top=content_top, bullet=bullet)

    # --- Two column layout ---
    elif layout_name == "two_column":
        if title:
            _set_title(slide, title, theme, font_size=Pt(32))

        if left_content:
            is_bullet = isinstance(left_content, list)
            _add_body_text(slide, left_content, theme,
                          left=MARGIN, top=CONTENT_TOP,
                          width=COLUMN_WIDTH, bullet=is_bullet)

        if right_content:
            is_bullet = isinstance(right_content, list)
            _add_body_text(slide, right_content, theme,
                          left=MARGIN + COLUMN_WIDTH + COLUMN_GAP,
                          top=CONTENT_TOP,
                          width=COLUMN_WIDTH, bullet=is_bullet)

    # --- Blank layout ---
    elif layout_name == "blank":
        if title:
            _set_title(slide, title, theme, font_size=Pt(28))

        if table_def:
            err = _add_table(slide, table_def, theme)
            if err:
                return err
        elif image:
            img_path = file_output_dir / image
            err = _add_image(slide, img_path, image_position, theme)
            if err:
                return err
        elif content:
            _add_body_text(slide, content, theme, bullet=bullet)

    # --- Speaker notes ---
    if notes:
        _add_notes(slide, notes)

    return None


def pptx_create_handler(server: FastMCP) -> None:
    """Register the pptx_create tool."""

    @server.tool()
    async def pptx_create(
        filename: str,
        slides: list[dict],
        theme: str = "default",
    ):
        """Create a PowerPoint (.pptx) presentation with text, images, and tables.

        Creates a valid .pptx file with multiple slides supporting common layouts.
        The file is saved to the output directory and returned as an embedded resource.

        **Preferred workflow: create small, then edit**

        Rather than building a complete presentation in one call, create a minimal
        version first (e.g. title slide + one content slide), then use ``pptx_edit``
        to add slides, fix text, insert tables, or adjust content. Use
        ``pptx_slide_image`` to visually verify how slides render. This iterative
        approach is more reliable and easier to debug.

        **Slide layouts**

        - ``"title"`` — Title slide with ``title`` and optional ``subtitle``
        - ``"title_content"`` — Title + body content (text, bullets, table, or image)
        - ``"two_column"`` — Title + ``left_content`` and ``right_content`` side by side
        - ``"section_header"`` — Section divider with ``title`` and ``subtitle``
        - ``"blank"`` — Empty canvas; add ``title``, ``content``, ``table``, or ``image``

        **Slide definition fields**

        - ``layout`` (str): One of the layouts above (default: ``"title_content"``)
        - ``title`` (str | None): Slide title
        - ``subtitle`` (str | None): Subtitle (for ``title`` / ``section_header`` layouts)
        - ``content`` (str | list[str] | None): Body text. A list renders as bullet points.
        - ``left_content`` (str | list[str] | None): Left column content (``two_column`` layout)
        - ``right_content`` (str | list[str] | None): Right column content (``two_column`` layout)
        - ``image`` (str | None): Filename of an image in the output directory
        - ``image_position`` (str): ``"left"``, ``"right"``, or ``"center"`` (default: ``"center"``)
        - ``table`` (dict | None): Table data with ``headers`` and ``rows``
        - ``notes`` (str | None): Speaker notes

        **Table definition**

        - ``headers`` (list[str]): Column headers
        - ``rows`` (list[list]): Data rows, each a list of values

        **Themes**

        - ``"default"`` — Standard PowerPoint theme
        - ``"dark"`` — Dark background with light text
        - ``"minimal"`` — Clean white with muted gray text

        The created file is returned as an MCP EmbeddedResource (base64 blob)
        so the LLM can read it directly, plus a TextContent summary.

        Args:
            filename: Output filename, must end with .pptx (e.g., "presentation.pptx")
            slides: List of slide definitions
            theme: Theme name (default: "default")

        Returns:
            MCP EmbeddedResource with base64-encoded .pptx file, plus a
            TextContent summary with slide details.
        """
        try:
            # --- Validate filename ---------------------------------------------
            err = _validate_filename(filename)
            if err:
                return [TextContent(type="text", text=err)]

            # --- Validate theme ------------------------------------------------
            if theme not in _VALID_THEMES:
                valid = ", ".join(sorted(_VALID_THEMES))
                return [TextContent(type="text",
                    text=f"Error: invalid theme '{theme}', must be {valid}")]

            # --- Validate slides -----------------------------------------------
            if not slides:
                return [TextContent(type="text", text="Error: at least one slide is required")]

            for i, slide_def in enumerate(slides):
                layout = slide_def.get("layout", "title_content")
                if layout not in _VALID_LAYOUTS:
                    valid = ", ".join(sorted(_VALID_LAYOUTS))
                    return [TextContent(type="text",
                        text=f"Error: slide {i + 1} has invalid layout '{layout}'")]

            # --- Create presentation -------------------------------------------
            prs = Presentation()
            prs.slide_width = SLIDE_WIDTH
            prs.slide_height = SLIDE_HEIGHT

            theme_colors = _THEMES[theme]
            file_output_dir = Path(settings.FILE_OUTPUT_DIR)

            for slide_def in slides:
                err = _build_slide(prs, slide_def, theme_colors, file_output_dir)
                if err:
                    return [TextContent(type="text", text=err)]

            # --- Save to buffer ------------------------------------------------
            output = BytesIO()
            prs.save(output)
            output.seek(0)
            pptx_bytes = output.getvalue()

            # --- Write to disk -------------------------------------------------
            file_output_dir.mkdir(parents=True, exist_ok=True)
            file_path = file_output_dir / filename
            file_path.write_bytes(pptx_bytes)

            # --- Build response ------------------------------------------------
            b64_content = base64.b64encode(pptx_bytes).decode("utf-8")
            resource_uri = f"file://{file_path}"

            embedded = EmbeddedResource(
                type="resource",
                resource={
                    "uri": resource_uri,
                    "mimeType": (
                        "application/vnd.openxmlformats-officedocument"
                        ".presentationml.presentation"),
                    "blob": b64_content,
                },
                annotations={"priority": 1},
            )

            # Build summary
            layout_counts = {}
            for slide_def in slides:
                layout = slide_def.get("layout", "title_content")
                layout_counts[layout] = layout_counts.get(layout, 0) + 1

            layout_summary = ", ".join(f"{v}x {k}" for k, v in sorted(layout_counts.items()))

            info_text = (
                f"PowerPoint presentation created successfully.\n"
                f"  Resource URI: {resource_uri}\n"
                f"  Download URL: {settings.FILE_BASE_URL}/files/{filename}\n"
                f"  Size: {len(pptx_bytes)} bytes\n"
                f"  Slides: {len(slides)}\n"
                f"  Layouts: {layout_summary}\n"
                f"  Theme: {theme}\n"
                f"\n"
                f"The file is embedded above as a base64 MCP resource. "
                f"You can also access it later via the resource URI or download URL."
            )

            return [embedded, TextContent(type="text", text=info_text)]

        except PermissionError as e:
            logger.error("Permission error creating file %s: %s", filename, str(e))
            msg = f"Error: permission denied writing to {settings.FILE_OUTPUT_DIR}"
            return [TextContent(type="text", text=f"{msg} - {str(e)}")]
        except OSError as e:
            logger.error("OS error creating file %s: %s", filename, str(e))
            return [TextContent(type="text", text=f"Error: failed to create file - {str(e)}")]
        except Exception as e:
            logger.error("Unexpected error creating file %s: %s", filename, str(e))
            return [TextContent(type="text", text=f"Error: {str(e)}")]

    logger.info("Registered pptx_create tool")
