"""PowerPoint (.pptx) edit tool for MCP server."""

import base64
import logging
from io import BytesIO
from pathlib import Path

from mcp.server import FastMCP
from mcp.types import EmbeddedResource, TextContent
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.config import settings

logger = logging.getLogger(__name__)

# Slide layout indices (matches pptx_create.py)
_LAYOUT_MAP = {
    "title": 0,
    "title_content": 1,
    "section_header": 2,
    "two_column": 3,
    "blank": 6,
}
_VALID_LAYOUTS = set(_LAYOUT_MAP.keys())
_VALID_THEMES = {"default", "dark", "minimal"}
_VALID_IMAGE_POSITIONS = {"left", "right", "center"}

# Theme colors (matches pptx_create.py)
_THEMES = {
    "default": {"bg": None, "title_color": None, "body_color": None, "accent": RGBColor(0x44, 0x72, 0xC4)},
    "dark": {"bg": RGBColor(0x1E, 0x1E, 0x1E), "title_color": RGBColor(0xFF, 0xFF, 0xFF), "body_color": RGBColor(0xCC, 0xCC, 0xCC), "accent": RGBColor(0x00, 0xB4, 0xD8)},
    "minimal": {"bg": RGBColor(0xFF, 0xFF, 0xFF), "title_color": RGBColor(0x33, 0x33, 0x33), "body_color": RGBColor(0x55, 0x55, 0x55), "accent": RGBColor(0x66, 0x66, 0x66)},
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.5)
CONTENT_TOP = Inches(1.8)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN * 2
COLUMN_GAP = Inches(0.4)
COLUMN_WIDTH = (CONTENT_WIDTH - COLUMN_GAP) / 2
TABLE_TOP = Inches(2.0)
TABLE_HEIGHT = Inches(4.5)
TABLE_COL_MIN = Inches(1.0)


def _validate_filename(filename: str) -> str | None:
    if not filename:
        return "Error: filename is required"
    if "/" in filename or "\\" in filename or ".." in filename:
        return "Error: filename must not contain path separators or '..'"
    if not filename.lower().endswith(".pptx"):
        return "Error: filename must end with .pptx"
    return None


def _set_run_color(run, color: RGBColor | None) -> None:
    if color is not None:
        run.font.color.rgb = color


def _add_title_text(slide, text: str, theme: dict,
                    left=None, top=None, width=None, height=None,
                    font_size=Pt(36), bold=True, alignment=PP_ALIGN.LEFT) -> None:
    if left is None:
        left = MARGIN
    if top is None:
        top = Inches(0.5)
    if width is None:
        width = CONTENT_WIDTH
    if height is None:
        height = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    _set_run_color(run, theme["title_color"])


def _add_body_text(slide, content, theme: dict,
                   left=None, top=None, width=None, height=None,
                   font_size=Pt(18), bullet=False,
                   alignment=PP_ALIGN.LEFT) -> None:
    if left is None:
        left = MARGIN
    if top is None:
        top = CONTENT_TOP
    if width is None:
        width = CONTENT_WIDTH
    if height is None:
        height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    items = content if isinstance(content, list) else [content]
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        if bullet:
            p.level = 0
            p.text = ""
            run = p.add_run()
            run.text = f"\u2022 {item}"
        else:
            p.text = ""
            run = p.add_run()
            run.text = item
        run.font.size = font_size
        _set_run_color(run, theme["body_color"])


def _add_table(slide, table_def: dict, theme: dict,
               left=None, top=None, width=None, height=None) -> str | None:
    headers = table_def.get("headers", [])
    rows = table_def.get("rows", [])
    if not headers:
        return "Error: table requires 'headers'"
    num_rows = 1 + len(rows)
    num_cols = len(headers)
    if left is None:
        left = MARGIN
    if top is None:
        top = TABLE_TOP
    if width is None:
        width = CONTENT_WIDTH
    if height is None:
        height = TABLE_HEIGHT
    min_total = TABLE_COL_MIN * num_cols
    if width < min_total:
        width = min_total
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    col_width = width // num_cols
    for i in range(num_cols):
        table.columns[i].width = col_width
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(header)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(14)
                if theme["title_color"]:
                    run.font.color.rgb = theme["title_color"]
                else:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.alignment = PP_ALIGN.CENTER
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = theme["accent"]
    for r_idx, row in enumerate(rows, 1):
        for c_idx, value in enumerate(row):
            if c_idx >= num_cols:
                break
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value) if value is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    _set_run_color(run, theme["body_color"])
    return None


def _add_image(slide, image_path: Path, position: str, theme: dict) -> str | None:
    if not image_path.exists():
        return f"Error: image file not found: {image_path.name}"
    img_width = Inches(5.0)
    img_height = Inches(4.0)
    if position == "left":
        left = MARGIN
        top = CONTENT_TOP
    elif position == "right":
        left = SLIDE_WIDTH - MARGIN - img_width
        top = CONTENT_TOP
    elif position == "center":
        left = (SLIDE_WIDTH - img_width) / 2
        top = (SLIDE_HEIGHT - img_height) / 2
    else:
        left = MARGIN
        top = CONTENT_TOP
    try:
        slide.shapes.add_picture(str(image_path), left, top, img_width, img_height)
    except Exception as e:
        return f"Error: failed to add image - {str(e)}"
    return None


def _build_slide(prs: "Presentation", slide_def: dict, theme: dict,
                 file_output_dir: Path) -> str | None:
    """Build a single slide from a slide definition. Returns error or None."""
    layout_name = slide_def.get("layout", "title_content")
    if layout_name not in _VALID_LAYOUTS:
        valid = ", ".join(sorted(_VALID_LAYOUTS))
        return f"Error: invalid layout '{layout_name}', must be {valid}"
    layout_idx = _LAYOUT_MAP[layout_name]
    if layout_idx >= len(prs.slide_layouts):
        layout_idx = len(prs.slide_layouts) - 1
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)

    title = slide_def.get("title")
    subtitle = slide_def.get("subtitle")
    content = slide_def.get("content")
    if content is None:
        content = slide_def.get("body")
    bullet = slide_def.get("bullet")
    if bullet is None:
        bullet = isinstance(content, list) if content else False
    notes = slide_def.get("notes")
    image = slide_def.get("image")
    image_position = slide_def.get("image_position", "center")
    table_def = slide_def.get("table")
    left_content = slide_def.get("left_content")
    right_content = slide_def.get("right_content")

    if layout_name == "title":
        if title:
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame:
                tf = title_shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.text = ""
                run = p.add_run()
                run.text = title
                run.font.size = Pt(44)
                run.font.bold = True
                _set_run_color(run, theme["title_color"])
            else:
                _add_title_text(slide, title, theme, font_size=Pt(44), alignment=PP_ALIGN.CENTER)
        if subtitle:
            _add_body_text(slide, subtitle, theme, top=Inches(4.0), font_size=Pt(24), alignment=PP_ALIGN.CENTER)

    elif layout_name == "section_header":
        if title:
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame:
                tf = title_shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.text = ""
                run = p.add_run()
                run.text = title
                run.font.size = Pt(40)
                run.font.bold = True
                _set_run_color(run, theme["title_color"])
            else:
                _add_title_text(slide, title, theme, font_size=Pt(40), alignment=PP_ALIGN.CENTER)
        if subtitle:
            _add_body_text(slide, subtitle, theme, top=Inches(4.2), font_size=Pt(22), alignment=PP_ALIGN.CENTER)

    elif layout_name == "title_content":
        if title:
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame:
                tf = title_shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = ""
                run = p.add_run()
                run.text = title
                run.font.size = Pt(32)
                run.font.bold = True
                _set_run_color(run, theme["title_color"])
            else:
                _add_title_text(slide, title, theme, font_size=Pt(32))
        if table_def:
            err = _add_table(slide, table_def, theme, top=CONTENT_TOP)
            if err:
                return err
        elif image:
            img_path = file_output_dir / image
            err = _add_image(slide, img_path, image_position, theme)
            if err:
                return err
            if content:
                if image_position in ("left", "right"):
                    text_left = (MARGIN if image_position == "right" else SLIDE_WIDTH - MARGIN - COLUMN_WIDTH)
                    _add_body_text(slide, content, theme, left=text_left, top=CONTENT_TOP, width=COLUMN_WIDTH)
                else:
                    _add_body_text(slide, content, theme, top=Inches(6.2), height=Inches(1.0))
        elif content:
            _add_body_text(slide, content, theme, top=CONTENT_TOP, bullet=bullet)

    elif layout_name == "two_column":
        if title:
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame:
                tf = title_shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = ""
                run = p.add_run()
                run.text = title
                run.font.size = Pt(32)
                run.font.bold = True
                _set_run_color(run, theme["title_color"])
            else:
                _add_title_text(slide, title, theme, font_size=Pt(32))
        if left_content:
            _add_body_text(slide, left_content, theme, left=MARGIN, top=CONTENT_TOP, width=COLUMN_WIDTH, bullet=isinstance(left_content, list))
        if right_content:
            _add_body_text(slide, right_content, theme, left=MARGIN + COLUMN_WIDTH + COLUMN_GAP, top=CONTENT_TOP, width=COLUMN_WIDTH, bullet=isinstance(right_content, list))

    elif layout_name == "blank":
        if title:
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame:
                tf = title_shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = ""
                run = p.add_run()
                run.text = title
                run.font.size = Pt(28)
                run.font.bold = True
                _set_run_color(run, theme["title_color"])
            else:
                _add_title_text(slide, title, theme, font_size=Pt(28))
        if table_def:
            err = _add_table(slide, table_def, theme)
            if err:
                return err
        elif image:
            img_path = file_output_dir / image
            err = _add_image(slide, img_path, image_position, theme)
            if err:
                return err
        elif content:
            _add_body_text(slide, content, theme, bullet=bullet)

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return None


def pptx_edit_handler(server: FastMCP) -> None:
    """Register the pptx_edit tool."""

    @server.tool()
    async def pptx_edit(
        filename: str,
        operations: list[dict],
        theme: str = "default",
    ):
        """Edit an existing PowerPoint (.pptx) presentation with targeted operations.

        Loads an existing presentation, applies one or more edit operations,
        and saves the result. This is the preferred way to incrementally
        build or modify presentations — create a small deck first with
        ``pptx_create``, then refine it with ``pptx_edit``.

        **Preferred workflow: create small, then edit**

        Rather than building a complete presentation in one ``pptx_create`` call,
        create a minimal version first (e.g. title slide + one content slide),
        then use ``pptx_edit`` to add slides, fix text, insert tables, or
        adjust content. This iterative approach is more reliable and easier
        to debug.

        **Operations**

        Each operation in ``operations`` is a dict with a ``type`` key and
        operation-specific parameters. Operations are applied in order.

        - ``"add_slide"`` — Add a new slide at the end or at a specific position
          - ``layout`` (str): Layout name (``"title"``, ``"title_content"``, ``"section_header"``, ``"two_column"``, ``"blank"``)
          - ``title`` (str, optional): Slide title
          - ``subtitle`` (str, optional): Subtitle (for ``title`` / ``section_header`` layouts)
          - ``content`` (str | list[str], optional): Body text; a list renders as bullets
          - ``left_content`` (str | list[str], optional): Left column (``two_column`` layout)
          - ``right_content`` (str | list[str], optional): Right column (``two_column`` layout)
          - ``table`` (dict, optional): Table with ``headers`` and ``rows``
          - ``image`` (str, optional): Filename of an image in the output directory
          - ``image_position`` (str): ``"left"``, ``"right"``, or ``"center"`` (default)
          - ``notes`` (str, optional): Speaker notes
          - ``at_index`` (int, optional): Insert at this 0-based position (default: append)

        - ``"update_slide_title"`` — Change a slide's title text
          - ``slide_index`` (int): 0-based slide index
          - ``title`` (str): New title text

        - ``"update_slide_content"`` — Replace body text on a slide
          - ``slide_index`` (int): 0-based slide index
          - ``content`` (str | list[str]): New body content; a list renders as bullets
          - ``bullet`` (bool, optional): Force bullet rendering

        - ``"update_slide_table"`` — Replace or add a table on a slide
          - ``slide_index`` (int): 0-based slide index
          - ``table`` (dict): Table with ``headers`` and ``rows``

        - ``"update_slide_notes"`` — Set speaker notes on a slide
          - ``slide_index`` (int): 0-based slide index
          - ``notes`` (str): Notes text

        - ``"delete_slide"`` — Remove a slide (must leave at least one slide)
          - ``slide_index`` (int): 0-based slide index to delete

        - ``"reorder_slide"`` — Move a slide to a new position
          - ``slide_index`` (int): Current 0-based index
          - ``to_index`` (int): New 0-based position

        - ``"add_slide_image"`` — Add an image to an existing slide
          - ``slide_index`` (int): 0-based slide index
          - ``image`` (str): Filename of an image in the output directory
          - ``position`` (str): ``"left"``, ``"right"``, or ``"center"`` (default)

        Args:
            filename: Existing .pptx file to edit
            operations: List of edit operations to apply in order
            theme: Theme name for new slides (default: "default")

        Returns:
            MCP EmbeddedResource with base64-encoded .pptx file, plus a
            TextContent summary of changes made.
        """
        try:
            # --- Validate filename ---------------------------------------------
            err = _validate_filename(filename)
            if err:
                return [TextContent(type="text", text=err)]

            # --- Validate theme ------------------------------------------------
            if theme not in _VALID_THEMES:
                valid = ", ".join(sorted(_VALID_THEMES))
                return [TextContent(type="text",
                    text=f"Error: invalid theme '{theme}', must be {valid}")]

            # --- Resolve file path ---------------------------------------------
            output_dir = Path(settings.FILE_OUTPUT_DIR)
            file_path = output_dir / filename

            if not file_path.exists():
                return [TextContent(
                    type="text",
                    text=f"Error: file '{filename}' not found in {settings.FILE_OUTPUT_DIR}",
                )]

            # --- Validate operations -------------------------------------------
            if not operations:
                return [TextContent(type="text", text="Error: at least one operation is required")]

            valid_types = {
                "add_slide", "update_slide_title", "update_slide_content",
                "update_slide_table", "update_slide_notes", "delete_slide",
                "reorder_slide", "add_slide_image",
            }
            for i, op in enumerate(operations):
                op_type = op.get("type")
                if op_type not in valid_types:
                    valid = ", ".join(sorted(valid_types))
                    return [TextContent(type="text",
                        text=f"Error: operation {i + 1} has invalid type '{op_type}', must be {valid}")]

            # --- Load presentation ---------------------------------------------
            prs = Presentation(str(file_path))
            theme_colors = _THEMES[theme]

            # --- Apply operations ----------------------------------------------
            changes = []

            for i, op in enumerate(operations):
                op_type = op["type"]

                # --- add_slide ---
                if op_type == "add_slide":
                    layout = op.get("layout", "title_content")
                    if layout not in _VALID_LAYOUTS:
                        valid = ", ".join(sorted(_VALID_LAYOUTS))
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} has invalid layout '{layout}', must be {valid}")]

                    at_index = op.get("at_index")
                    if at_index is not None:
                        if at_index < 0 or at_index > len(prs.slides):
                            return [TextContent(type="text",
                                text=f"Error: operation {i + 1} at_index {at_index} out of range (0-{len(prs.slides)})")]

                    # Build slide at the end, then move if needed
                    slide_def = {
                        "layout": layout,
                        "title": op.get("title"),
                        "subtitle": op.get("subtitle"),
                        "content": op.get("content"),
                        "body": op.get("body"),
                        "bullet": op.get("bullet"),
                        "left_content": op.get("left_content"),
                        "right_content": op.get("right_content"),
                        "table": op.get("table"),
                        "image": op.get("image"),
                        "image_position": op.get("image_position", "center"),
                        "notes": op.get("notes"),
                    }
                    err = _build_slide(prs, slide_def, theme_colors, output_dir)
                    if err:
                        return [TextContent(type="text", text=err)]

                    if at_index is not None and at_index < len(prs.slides) - 1:
                        # Move the newly added slide to the desired position
                        rId = prs.slides._sldIdLst[-1].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                        sldId = prs.slides._sldIdLst[-1]
                        prs.slides._sldIdLst.remove(sldId)
                        prs.slides._sldIdLst.insert(at_index, sldId)

                    changes.append(f"Added slide with layout '{layout}'" + (f" at index {at_index}" if at_index is not None else ""))

                # --- update_slide_title ---
                elif op_type == "update_slide_title":
                    slide_index = op.get("slide_index")
                    title = op.get("title")
                    if slide_index is None or title is None:
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} (update_slide_title) requires 'slide_index' and 'title'")]
                    if slide_index < 0 or slide_index >= len(prs.slides):
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} slide_index {slide_index} out of range (0-{len(prs.slides) - 1})")]
                    slide = prs.slides[slide_index]
                    title_shape = slide.shapes.title
                    if title_shape and title_shape.has_text_frame:
                        tf = title_shape.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.text = ""
                        run = p.add_run()
                        run.text = title
                        run.font.size = Pt(32)
                        run.font.bold = True
                        _set_run_color(run, theme_colors["title_color"])
                    else:
                        _add_title_text(slide, title, theme_colors, font_size=Pt(32))
                    changes.append(f"Updated title on slide {slide_index + 1}")

                # --- update_slide_content ---
                elif op_type == "update_slide_content":
                    slide_index = op.get("slide_index")
                    content = op.get("content")
                    if slide_index is None or content is None:
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} (update_slide_content) requires 'slide_index' and 'content'")]
                    if slide_index < 0 or slide_index >= len(prs.slides):
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} slide_index {slide_index} out of range (0-{len(prs.slides) - 1})")]
                    slide = prs.slides[slide_index]
                    bullet = op.get("bullet")
                    if bullet is None:
                        bullet = isinstance(content, list)
                    _add_body_text(slide, content, theme_colors, top=CONTENT_TOP, bullet=bullet)
                    changes.append(f"Updated content on slide {slide_index + 1}")

                # --- update_slide_table ---
                elif op_type == "update_slide_table":
                    slide_index = op.get("slide_index")
                    table_def = op.get("table")
                    if slide_index is None or not table_def:
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} (update_slide_table) requires 'slide_index' and 'table'")]
                    if slide_index < 0 or slide_index >= len(prs.slides):
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} slide_index {slide_index} out of range (0-{len(prs.slides) - 1})")]
                    slide = prs.slides[slide_index]
                    err = _add_table(slide, table_def, theme_colors, top=CONTENT_TOP)
                    if err:
                        return [TextContent(type="text", text=err)]
                    changes.append(f"Added/updated table on slide {slide_index + 1}")

                # --- update_slide_notes ---
                elif op_type == "update_slide_notes":
                    slide_index = op.get("slide_index")
                    notes = op.get("notes")
                    if slide_index is None or notes is None:
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} (update_slide_notes) requires 'slide_index' and 'notes'")]
                    if slide_index < 0 or slide_index >= len(prs.slides):
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} slide_index {slide_index} out of range (0-{len(prs.slides) - 1})")]
                    slide = prs.slides[slide_index]
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = notes
                    changes.append(f"Updated notes on slide {slide_index + 1}")

                # --- delete_slide ---
                elif op_type == "delete_slide":
                    slide_index = op.get("slide_index")
                    if slide_index is None:
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} (delete_slide) requires 'slide_index'")]
                    if slide_index < 0 or slide_index >= len(prs.slides):
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} slide_index {slide_index} out of range (0-{len(prs.slides) - 1})")]
                    if len(prs.slides) <= 1:
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} cannot delete the last slide")]
                    # Remove the slide's XML element
                    rId = prs.slides._sldIdLst[slide_index].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[slide_index]
                    changes.append(f"Deleted slide at index {slide_index + 1}")

                # --- reorder_slide ---
                elif op_type == "reorder_slide":
                    slide_index = op.get("slide_index")
                    to_index = op.get("to_index")
                    if slide_index is None or to_index is None:
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} (reorder_slide) requires 'slide_index' and 'to_index'")]
                    if slide_index < 0 or slide_index >= len(prs.slides):
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} slide_index {slide_index} out of range (0-{len(prs.slides) - 1})")]
                    if to_index < 0 or to_index >= len(prs.slides):
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} to_index {to_index} out of range (0-{len(prs.slides) - 1})")]
                    if slide_index == to_index:
                        changes.append(f"Slide {slide_index + 1} already at position {to_index + 1}")
                        continue
                    sldId = prs.slides._sldIdLst[slide_index]
                    prs.slides._sldIdLst.remove(sldId)
                    prs.slides._sldIdLst.insert(to_index, sldId)
                    changes.append(f"Moved slide {slide_index + 1} to position {to_index + 1}")

                # --- add_slide_image ---
                elif op_type == "add_slide_image":
                    slide_index = op.get("slide_index")
                    image = op.get("image")
                    if slide_index is None or not image:
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} (add_slide_image) requires 'slide_index' and 'image'")]
                    if slide_index < 0 or slide_index >= len(prs.slides):
                        return [TextContent(type="text",
                            text=f"Error: operation {i + 1} slide_index {slide_index} out of range (0-{len(prs.slides) - 1})")]
                    position = op.get("position", "center")
                    if position not in _VALID_IMAGE_POSITIONS:
                        valid = ", ".join(sorted(_VALID_IMAGE_POSITIONS))
                        return [TextContent(type="text",
                            text=f"Error: invalid position '{position}', must be {valid}")]
                    slide = prs.slides[slide_index]
                    img_path = output_dir / image
                    err = _add_image(slide, img_path, position, theme_colors)
                    if err:
                        return [TextContent(type="text", text=err)]
                    changes.append(f"Added image '{image}' to slide {slide_index + 1}")

            # --- Save to buffer ------------------------------------------------
            output = BytesIO()
            prs.save(output)
            output.seek(0)
            pptx_bytes = output.getvalue()

            # --- Write to disk -------------------------------------------------
            file_path.write_bytes(pptx_bytes)

            # --- Build response ------------------------------------------------
            b64_content = base64.b64encode(pptx_bytes).decode("utf-8")
            resource_uri = f"file://{file_path}"

            embedded = EmbeddedResource(
                type="resource",
                resource={
                    "uri": resource_uri,
                    "mimeType": (
                        "application/vnd.openxmlformats-officedocument"
                        ".presentationml.presentation"),
                    "blob": b64_content,
                },
                annotations={"priority": 1},
            )

            info_text = (
                f"PowerPoint presentation edited successfully.\n"
                f"  Resource URI: {resource_uri}\n"
                f"  Size: {len(pptx_bytes)} bytes\n"
                f"  Slides: {len(prs.slides)}\n"
                f"  Operations applied: {len(operations)}\n"
                + "\n".join(f"  - {c}" for c in changes)
                + "\n\n"
                "The file is embedded above as a base64 MCP resource. "
                "You can also access it later via the resource URI."
            )

            return [embedded, TextContent(type="text", text=info_text)]

        except PermissionError as e:
            logger.error("Permission error editing file %s: %s", filename, str(e))
            msg = f"Error: permission denied writing to {settings.FILE_OUTPUT_DIR}"
            return [TextContent(type="text", text=f"{msg} - {str(e)}")]
        except OSError as e:
            logger.error("OS error editing file %s: %s", filename, str(e))
            return [TextContent(type="text", text=f"Error: failed to edit file - {str(e)}")]
        except Exception as e:
            logger.error("Unexpected error editing file %s: %s", filename, str(e))
            return [TextContent(type="text", text=f"Error: {str(e)}")]

    logger.info("Registered pptx_edit tool")
