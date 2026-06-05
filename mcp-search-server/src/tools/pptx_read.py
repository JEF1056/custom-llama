"""PowerPoint (.pptx) read tool for MCP server."""

import json
import logging
from pathlib import Path

from mcp.server import FastMCP
from mcp.types import TextContent
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.config import settings

logger = logging.getLogger(__name__)

# Layout index to name mapping (matches pptx_create.py)
_LAYOUT_NAMES = {
    0: "title",
    1: "title_content",
    2: "section_header",
    3: "two_column",
    6: "blank",
}


def _validate_filename(filename: str) -> str | None:
    """Validate filename for path traversal. Returns error message or None."""
    if not filename:
        return "Error: filename is required"
    if "/" in filename or "\\" in filename or ".." in filename:
        return "Error: filename must not contain path separators or '..'"
    if not filename.lower().endswith(".pptx"):
        return "Error: filename must end with .pptx"
    return None


def _extract_text_from_shape(shape) -> str:
    """Extract all text from a shape that has a text frame."""
    if not shape.has_text_frame:
        return ""
    parts = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    return "\n".join(parts)


def _extract_table_data(shape) -> dict:
    """Extract table data from a table shape."""
    if not shape.has_table:
        return {}
    table = shape.table
    headers = []
    rows = []
    for row_idx, row in enumerate(table.rows):
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text.strip())
        if row_idx == 0:
            headers = row_data
        else:
            rows.append(row_data)
    return {"headers": headers, "rows": rows}


def _describe_shape(shape) -> dict:
    """Describe a shape's type and content for the read output."""
    info = {"shape_id": shape.shape_id}

    if shape.has_text_frame:
        text = _extract_text_from_shape(shape)
        if text:
            info["text"] = text

    if shape.has_table:
        info["type"] = "table"
        info["table"] = _extract_table_data(shape)
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        info["type"] = "image"
        info["image"] = {
            "width": shape.width,
            "height": shape.height,
            "left": shape.left,
            "top": shape.top,
        }
    elif shape.has_text_frame:
        info["type"] = "text"
    else:
        info["type"] = "other"

    return info


def pptx_read_handler(server: FastMCP) -> None:
    """Register the pptx_read tool."""

    @server.tool()
    async def pptx_read(
        filename: str,
        slide_index: int | None = None,
    ):
        """Read a .pptx and return structured slide data (layout, title, shapes, notes).

        slide_index: read one slide (0-based); omit for all slides.
        """
        try:
            # --- Validate filename ---------------------------------------------
            err = _validate_filename(filename)
            if err:
                return [TextContent(type="text", text=err)]

            # --- Resolve file path ---------------------------------------------
            output_dir = Path(settings.FILE_OUTPUT_DIR)
            file_path = output_dir / filename

            if not file_path.exists():
                return [TextContent(
                    type="text",
                    text=f"Error: file '{filename}' not found in {settings.FILE_OUTPUT_DIR}",
                )]

            # --- Load presentation ---------------------------------------------
            prs = Presentation(str(file_path))

            # --- Determine which slides to read --------------------------------
            slide_indices = range(len(prs.slides))
            if slide_index is not None:
                if slide_index < 0 or slide_index >= len(prs.slides):
                    return [TextContent(
                        type="text",
                        text=f"Error: slide index {slide_index} out of range (0-{len(prs.slides) - 1})",
                    )]
                slide_indices = [slide_index]

            # --- Read each slide -----------------------------------------------
            slides_data = []
            for idx in slide_indices:
                slide = prs.slides[idx]

                # Determine layout name
                layout_name = "unknown"
                slide_layout = slide.slide_layout
                for layout_idx, name in _LAYOUT_NAMES.items():
                    if layout_idx < len(prs.slide_layouts):
                        if prs.slide_layouts[layout_idx] == slide_layout:
                            layout_name = name
                            break

                # Extract title from title placeholder
                title_text = ""
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    title_text = slide.shapes.title.text_frame.text.strip()

                # Describe all shapes
                shapes = []
                for shape in slide.shapes:
                    shape_info = _describe_shape(shape)
                    shapes.append(shape_info)

                # Extract speaker notes
                notes_text = ""
                if slide.has_notes_slide and slide.notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame:
                        notes_text = notes_frame.text.strip()

                slides_data.append({
                    "index": idx,
                    "layout": layout_name,
                    "title": title_text,
                    "shapes": shapes,
                    "notes": notes_text,
                })

            # --- Build response ------------------------------------------------
            result = {
                "status": "success",
                "filename": filename,
                "slide_count": len(prs.slides),
                "slides": slides_data,
            }

            return [TextContent(type="text", text=json.dumps(result, indent=2))]

        except PermissionError as e:
            logger.error("Permission error reading file %s: %s", filename, str(e))
            msg = f"Error: permission denied reading {filename}"
            return [TextContent(type="text", text=f"{msg} - {str(e)}")]
        except Exception as e:
            logger.error("Unexpected error reading file %s: %s", filename, str(e))
            return [TextContent(type="text", text=f"Error: failed to read '{filename}' - {str(e)}")]

    logger.info("Registered pptx_read tool")
