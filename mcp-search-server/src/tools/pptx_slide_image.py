"""PowerPoint (.pptx) slide image tool for MCP server."""

import base64
import logging
import os
import shutil
import subprocess
import tempfile
from pathlib import Path

from mcp.server import FastMCP
from mcp.types import ImageContent, TextContent

from src.config import settings

logger = logging.getLogger(__name__)


def _validate_filename(filename: str) -> str | None:
    if not filename:
        return "Error: filename is required"
    if "/" in filename or "\\" in filename or ".." in filename:
        return "Error: filename must not contain path separators or '..'"
    if not filename.lower().endswith(".pptx"):
        return "Error: filename must end with .pptx"
    return None


def _find_libreoffice() -> str | None:
    """Find the LibreOffice executable."""
    candidates = [
        "libreoffice",
        "soffice",
        "/usr/bin/libreoffice",
        "/usr/bin/soffice",
        "/usr/lib/libreoffice/program/soffice",
        "/opt/libreoffice/program/soffice",
        "/snap/bin/libreoffice",
    ]
    for candidate in candidates:
        if shutil.which(candidate):
            return candidate
    return None


def _convert_slide_to_image(file_path: Path, slide_index: int, output_dir: str) -> str | None:
    """Convert a specific slide to a PNG image using LibreOffice.

    LibreOffice converts the entire presentation to one PNG per slide,
    named like 'filename-1.png', 'filename-2.png', etc.

    Returns the path to the image file, or None if conversion failed.
    """
    lo_path = _find_libreoffice()
    if not lo_path:
        return None

    try:
        result = subprocess.run(
            [
                lo_path,
                "--headless",
                "--convert-to", "png",
                "--outdir", output_dir,
                str(file_path),
            ],
            capture_output=True,
            text=True,
            timeout=60,
            env={**os.environ, "HOME": output_dir},
        )
        if result.returncode != 0:
            logger.warning("LibreOffice conversion failed: %s", result.stderr)
            return None

        # LibreOffice names output as 'basename-1.png', 'basename-2.png', etc.
        # The slide_index is 0-based, so we need slide_index + 1
        stem = file_path.stem
        image_name = f"{stem}-{slide_index + 1}.png"
        image_path = Path(output_dir) / image_name

        if image_path.exists():
            return str(image_path)

        # Try alternative naming: some versions use 0-based indexing
        image_name = f"{stem}-{slide_index}.png"
        image_path = Path(output_dir) / image_name
        if image_path.exists():
            return str(image_path)

        # List what was actually created for debugging
        created = list(Path(output_dir).glob("*.png"))
        if created:
            logger.warning("Expected %s but found: %s", image_name, [f.name for f in created])
            # Return the closest match
            return str(created[0])

        return None

    except subprocess.TimeoutExpired:
        logger.warning("LibreOffice conversion timed out")
        return None
    except Exception as e:
        logger.warning("LibreOffice conversion error: %s", str(e))
        return None


def pptx_slide_image_handler(server: FastMCP) -> None:
    """Register the pptx_slide_image tool."""

    @server.tool()
    async def pptx_slide_image(
        filename: str,
        slide_index: int,
    ):
        """Render a specific slide from a PowerPoint (.pptx) as an image.

        Converts the requested slide to a PNG image using LibreOffice
        (if available) and returns it as an MCP image so the LLM can
        visually inspect the slide layout, formatting, and content.

        **Preferred workflow: create small, then edit**

        After creating or editing a presentation, use this tool to verify
        how slides actually render. This is especially useful after
        ``pptx_edit`` operations to confirm changes look correct.

        **Parameters**

        - ``filename`` (str): Name of the .pptx file (e.g., "presentation.pptx").
          Must not contain path separators or ``..``.
        - ``slide_index`` (int): 0-based index of the slide to render.

        **Returns**

        - On success: MCP ImageContent with the slide as a PNG image
        - On failure (LibreOffice not available): TextContent with structured
          slide data from ``pptx_read`` as a fallback

        Args:
            filename: Name of the .pptx file
            slide_index: 0-based index of the slide to render

        Returns:
            MCP ImageContent with the slide PNG, or TextContent with fallback data.
        """
        try:
            # --- Validate filename ---------------------------------------------
            err = _validate_filename(filename)
            if err:
                return [TextContent(type="text", text=err)]

            # --- Resolve file path ---------------------------------------------
            output_dir = Path(settings.FILE_OUTPUT_DIR)
            file_path = output_dir / filename

            if not file_path.exists():
                return [TextContent(
                    type="text",
                    text=f"Error: file '{filename}' not found in {settings.FILE_OUTPUT_DIR}",
                )]

            # --- Validate slide index ------------------------------------------
            from pptx import Presentation
            prs = Presentation(str(file_path))
            if slide_index < 0 or slide_index >= len(prs.slides):
                return [TextContent(
                    type="text",
                    text=f"Error: slide index {slide_index} out of range (0-{len(prs.slides) - 1}). Presentation has {len(prs.slides)} slide(s).",
                )]

            # --- Try LibreOffice conversion ------------------------------------
            with tempfile.TemporaryDirectory() as tmpdir:
                image_path = _convert_slide_to_image(file_path, slide_index, tmpdir)

                if image_path and Path(image_path).exists():
                    # Read and encode the image
                    img_bytes = Path(image_path).read_bytes()
                    b64_content = base64.b64encode(img_bytes).decode("utf-8")

                    return [
                        ImageContent(
                            type="image",
                            data=b64_content,
                            mimeType="image/png",
                        ),
                        TextContent(
                            type="text",
                            text=f"Slide {slide_index + 1} of {len(prs.slides)} rendered as PNG ({len(img_bytes)} bytes).",
                        ),
                    ]

            # --- Fallback: return structured slide data ------------------------
            slide = prs.slides[slide_index]
            title = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text.strip()

            shapes_info = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        shapes_info.append(f"  - Text: {text[:100]}")
                if shape.has_table:
                    shapes_info.append(f"  - Table: {len(shape.table.rows)} rows x {len(shape.table.columns)} cols")
                if hasattr(shape, 'image'):
                    shapes_info.append(f"  - Image: {shape.width}x{shape.height}")

            notes = ""
            if slide.has_notes_slide and slide.notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            fallback = (
                f"LibreOffice is not available for slide rendering.\n"
                f"Here is the structured data for slide {slide_index + 1}:\n\n"
                f"  Title: {title or '(none)'}\n"
                f"  Shapes:\n" + "\n".join(shapes_info) + "\n"
                f"  Notes: {notes or '(none)'}\n\n"
                f"Install LibreOffice in the container to enable visual slide rendering."
            )
            return [TextContent(type="text", text=fallback)]

        except PermissionError as e:
            logger.error("Permission error reading file %s: %s", filename, str(e))
            msg = f"Error: permission denied reading {filename}"
            return [TextContent(type="text", text=f"{msg} - {str(e)}")]
        except Exception as e:
            logger.error("Unexpected error rendering slide from %s: %s", filename, str(e))
            return [TextContent(type="text", text=f"Error: failed to render slide - {str(e)}")]

    logger.info("Registered pptx_slide_image tool")
